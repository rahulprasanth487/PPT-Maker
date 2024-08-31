from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.text.text import TextFrame


# --- SLIDE CONTENT FROM EXTRACTED DATA ---
slide_titles = [
    "SLICE",
    "About Slice",
    "Startup Story & Target Market",
    "Mission",
    "Business Model (Pre-RBI Change)",
    "Change in Business Model (Post July 20, 2022)",
    "Growth & Revenue Highlights",
    "Financial Performance FY22-FY23",
    "Expense to Revenue Ratio Analysis",
    "EBITDA Margin & ROCE",
    "Other Features"
]

slide_content = [
    # Slide 1: SLICE (Title Slide)
    "Credit Card Startup | Hassle-Free Payments | Converted to EMI", 

    # Slide 2: About Slice
    "- Credit startup offering hassle-free payment cards converted to EMI.\n"
    "- Slice Super Cards: Zero-fee cards for online and offline transactions.\n"
    "- 2% cashback on every transaction using Slice Credit Card.\n"
    "- Credit limit: ₹2,000 to ₹10 Lakhs\n"
    "- No interest if the amount is repaid within 3 EMIs.\n"
    "- In-app spending insights: track patterns, repayment reminders, and more.", 

    # Slide 3: Startup Story & Target Market
    "- Founded to address the hassles associated with traditional credit cards, particularly for middle-class families in Tier 2 & Tier 3 cities.\n"
    "- Target Market: Students and early-age business professionals ineligible for traditional credit cards.\n"
    "- Average age of Slice customers: 22 years.", 

    # Slide 4: Mission
    "- Provide accessible credit options for students and young professionals to purchase essential items like laptops and mobile phones.\n"
    "- Offer a more flexible and transparent repayment system through monthly installments.",

    # Slide 5: Business Model (Pre-RBI Change)
    "- Subvention Income: Partnerships with merchants like Amazon and Flipkart for no-cost EMIs.\n"
    "- Interchange Fees: Percentage-based fees on transactions from merchants.\n"
    "- Interest Income from EMIs.\n"
    "- Additional Fees: Late payment, foreign transaction, service charges, etc.\n"
    "- Targeted Approach: Focus on young adults and millennials.\n"
    "- Transparent Pricing: Building trust through clear fee structures.\n"
    "- Data-Driven Insights: Provide users with spending patterns and recommendations for EMI conversion.",

    # Slide 6: Change in Business Model (Post July 20, 2022)
    "- RBI Announcement Impact: Transitioned from offering credit cards to 'classic term loans'.\n"
    "- Purchase Power: Creditworthiness evaluated on a transaction-by-transaction basis, using factors like payment history and financial situation.", 

    # Slide 7: Funding Rounds (Consider a table format if possible)
    # ... (Add funding data here - I can't format this well as text) 

    # Slide 8: Growth & Revenue Highlights
    "- Crossed one million app transactions within 5 months of launching physical cards in May 2019.\n"
    "- Pivoted from BNPL to card products in 2019 and now includes UPI products.\n"
    "- Achieved unicorn status in 2021 with a valuation of $1.8 billion as of March 2023.",

    # Slide 9: Financial Performance FY22-FY23 (Use charts from document)
    "- Operating Revenue: Significant growth (refer to chart data).\n"
    "- Net Loss: Increased due to factors like advertising expenses (refer to chart data).",

    # Slide 10: Expense to Revenue Ratio Analysis
    "- FY21: Expense to Revenue Ratio ~ 248.46% (indicating high expenses).\n"
    "- FY23: Expense to Revenue Ratio ~ 150.21% (improvement, but still elevated).",

    # Slide 11: EBITDA Margin & ROCE
    "- EBITDA margin improved in FY22, indicating better cost management.\n"
    "- ROCE remained negative but showed improvement.",

    # Slide 12: Other Features
    "- Auto reload of money into the wallet.\n"
    "- Seamless transactions with auto-reload functionality.\n"
    "- Slice Card UPI."
]


font_family = "Arial"  # Or a close substitute
heading_font_size = Pt(26)
subheading_font_size = Pt(20) 
body_font_size = Pt(15)
wellsfargo_red = (235, 23, 30) # RGB approximation, needs confirmation
dark_gray = (51, 51, 51)  # Approximate gray for text

# --- CREATE PRESENTATION ---
prs = Presentation()
max_chars_per_line = 100


for i, title in enumerate(slide_titles):
    slide_layout = prs.slide_layouts[6]  # Blank layout for more control
    slide = prs.slides.add_slide(slide_layout)


    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
    title_frame = title_box.text_frame
    title_paragraph = title_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = heading_font_size
    title_paragraph.font.name = font_family
    title_paragraph.font.color.rgb = RGBColor(*wellsfargo_red)

    # Add Content (adapt as needed for bullet points, etc.)
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(10))
    content_frame = content_box.text_frame

    for paragraph_text in slide_content[i].split('\n'):
        if i == 7:
            words = paragraph_text.split()
            current_line = ""
            for word in words:
                if len(current_line) + len(word) + 1 <= 46:
                    current_line += word + " "
                else:
                    p = content_frame.add_paragraph() 
                    p.text = current_line
                    current_line = " "+word + " "
            p = content_frame.add_paragraph()  
            p.text = current_line
            slide.shapes.add_picture("images/image4.png", Inches(5), Inches(1.5), width=Inches(5))

        elif i == 8: 
            words = paragraph_text.split()
            current_line = ""
            for word in words:
                if len(current_line) + len(word) + 1 <= 46:
                    current_line += word + " "
                else:
                    p = content_frame.add_paragraph() 
                    p.text = current_line
                    current_line = " "+word + " "
            p = content_frame.add_paragraph()  
            p.text = current_line
            slide.shapes.add_picture("images/image5.png", Inches(5), Inches(1.5), width=Inches(5))

        elif i == 9:
            words = paragraph_text.split()
            current_line = ""
            for word in words:
                if len(current_line) + len(word) + 1 <= 46:
                    current_line += word + " "
                else:
                    p = content_frame.add_paragraph() 
                    p.text = current_line
                    current_line = " "+word + " "
            p = content_frame.add_paragraph()  
            p.text = current_line
            slide.shapes.add_picture("images/image6.png", Inches(5), Inches(1.5), width=Inches(5))

        else:
            words = paragraph_text.split()
            current_line = ""
            for word in words:
                if len(current_line) + len(word) + 1 <= max_chars_per_line:
                    current_line += word + " "
                else:
                    p = content_frame.add_paragraph() 
                    p.text = current_line
                    current_line = " "+word + " "
            p = content_frame.add_paragraph()  
            p.text = current_line

        # --- Formatting for EACH added paragraph ---
        for paragraph in content_frame.paragraphs: # Loop and format
            paragraph.font.size = body_font_size
            paragraph.font.name = font_family
            paragraph.font.color.rgb = RGBColor(*dark_gray)
            paragraph.space_after = Pt(10) 

        


# --- SAVE PRESENTATION ---
prs.save("slice_presentation_wellsfargo_style.pptx") 